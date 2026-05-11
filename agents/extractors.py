import importlib
import json
import mimetypes
import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any
from urllib.parse import unquote, urlparse
from urllib.request import urlopen

IMAGE_TYPES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"}
TEXT_TYPES = {".txt", ".md", ".json", ".xml", ".html", ".csv", ".yaml", ".yml"}
PDF_TYPES = {".pdf"}
DOCX_TYPES = {".docx"}
XLSX_TYPES = {".xlsx", ".xls"}
PPTX_TYPES = {".pptx"}
AUDIO_TYPES = {".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg"}
VIDEO_TYPES = {".mp4", ".mov", ".avi", ".mkv", ".webm"}

_AUDIO_MODEL: Any | None = None
_AUDIO_BACKEND: str | None = None


def is_url(value: str) -> bool:
    parsed = urlparse(value)
    return parsed.scheme in {"http", "https"} and bool(parsed.netloc)


def get_path_suffix(asset_path: str) -> str:
    if is_url(asset_path):
        return Path(unquote(urlparse(asset_path).path)).suffix.lower()
    return Path(asset_path).suffix.lower()


def infer_output_stem(asset_path: str) -> str:
    if is_url(asset_path):
        parsed = urlparse(asset_path)
        candidate = Path(unquote(parsed.path)).stem or parsed.netloc or "remote-asset"
        normalized = "".join(
            char if char.isalnum() or char in {"-", "_"} else "-" for char in candidate
        ).strip("-_")
        return normalized or "remote-asset"

    return Path(asset_path).stem


def infer_download_suffix(asset_path: str, default_suffix: str = "") -> str:
    suffix = get_path_suffix(asset_path)
    return suffix or default_suffix


def download_binary_url(asset_url: str, default_suffix: str = "") -> tuple[str, str]:
    download_dir = tempfile.mkdtemp(prefix="ocr-agent-remote-")
    suffix = infer_download_suffix(asset_url, default_suffix)
    filename = f"{infer_output_stem(asset_url)}{suffix}"
    local_path = Path(download_dir) / filename

    with urlopen(asset_url) as response, local_path.open("wb") as file_handle:
        file_handle.write(response.read())

    return str(local_path), download_dir


def resolve_local_file(asset_path: str, label: str) -> str:
    resolved_path = Path(asset_path).expanduser().resolve()
    if not resolved_path.exists():
        raise FileNotFoundError(f"{label} file does not exist: {resolved_path}")
    if not resolved_path.is_file():
        raise ValueError(f"{label} path is not a file: {resolved_path}")
    if resolved_path.stat().st_size == 0:
        raise ValueError(f"{label} file is empty: {resolved_path}")
    return str(resolved_path)


def resolve_remote_or_local_file(
    asset_path: str,
    *,
    label: str,
    default_suffix: str = "",
) -> tuple[str, str | None]:
    if is_url(asset_path):
        return download_binary_url(asset_path, default_suffix=default_suffix)
    return resolve_local_file(asset_path, label), None


def cleanup_temporary_resource(
    file_path: str | None = None,
    directory: str | None = None,
):
    if file_path and os.path.exists(file_path):
        os.unlink(file_path)
    if directory:
        shutil.rmtree(directory, ignore_errors=True)


def classify_image_text_type(image_path: str) -> str:
    from agents.vision_llm import classify_image_text_type_with_llm

    local_image_path, cleanup_dir = resolve_remote_or_local_file(
        image_path,
        label="Image",
        default_suffix=".png",
    )
    try:
        return classify_image_text_type_with_llm(local_image_path)
    finally:
        cleanup_temporary_resource(directory=cleanup_dir)


def extract_handwritten_image(asset_path: str) -> str:
    from agents.vision_llm import extract_handwritten_image_with_llm

    local_image_path, cleanup_dir = resolve_remote_or_local_file(
        asset_path,
        label="Image",
        default_suffix=".png",
    )
    try:
        return extract_handwritten_image_with_llm(local_image_path)
    finally:
        cleanup_temporary_resource(directory=cleanup_dir)


def detect_file_type(asset_path: str) -> str:
    suffix = get_path_suffix(asset_path)

    if "youtube.com" in asset_path or "youtu.be" in asset_path:
        return "video"
    if suffix in IMAGE_TYPES:
        return "image"
    if suffix in PDF_TYPES:
        return "pdf"
    if suffix in DOCX_TYPES:
        return "docx"
    if suffix in XLSX_TYPES:
        return "excel"
    if suffix in PPTX_TYPES:
        return "pptx"
    if suffix in AUDIO_TYPES:
        return "audio"
    if suffix in VIDEO_TYPES:
        return "video"
    if suffix in TEXT_TYPES:
        return "text"

    mime_type, _ = mimetypes.guess_type(asset_path)
    if mime_type:
        if mime_type.startswith("image/"):
            return "image"
        if mime_type.startswith("audio/"):
            return "audio"
        if mime_type.startswith("video/"):
            return "video"
        if mime_type == "application/pdf":
            return "pdf"

    return mime_type or "unknown"


def extract_text_file(asset_path: str) -> str:
    import pandas as pd

    path = Path(asset_path)

    if path.suffix.lower() == ".csv":
        return pd.read_csv(path).to_string(index=False)

    if path.suffix.lower() == ".json":
        data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
        return json.dumps(data, indent=2, ensure_ascii=False)

    return path.read_text(encoding="utf-8", errors="ignore")


def extract_pdf(asset_path: str) -> str:
    import fitz
    import pdfplumber
    import pytesseract
    from PIL import Image

    local_pdf_path, cleanup_dir = resolve_remote_or_local_file(
        asset_path,
        label="PDF",
        default_suffix=".pdf",
    )
    try:
        text_parts: list[str] = []

        with pdfplumber.open(local_pdf_path) as pdf:
            for page_number, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    text_parts.append(f"\n--- PAGE {page_number} ---\n{text}")

        if "".join(text_parts).strip():
            return "\n".join(text_parts)

        doc = fitz.open(local_pdf_path)
        for page_number, page in enumerate(doc, start=1):
            pix = page.get_pixmap(dpi=300)

            with tempfile.NamedTemporaryFile(suffix=".png", delete=True) as tmp:
                pix.save(tmp.name)
                image = Image.open(tmp.name)
                text = pytesseract.image_to_string(image)

            text_parts.append(f"\n--- PAGE {page_number} OCR ---\n{text}")

        return "\n".join(text_parts)
    finally:
        cleanup_temporary_resource(directory=cleanup_dir)


def extract_printed_image(asset_path: str) -> str:
    import pytesseract
    from PIL import Image

    local_image_path, cleanup_dir = resolve_remote_or_local_file(
        asset_path,
        label="Image",
        default_suffix=".png",
    )
    try:
        return pytesseract.image_to_string(Image.open(local_image_path))
    finally:
        cleanup_temporary_resource(directory=cleanup_dir)


def extract_pdf_as_images(asset_path: str) -> str:
    import fitz

    local_pdf_path, cleanup_dir = resolve_remote_or_local_file(
        asset_path,
        label="PDF",
        default_suffix=".pdf",
    )
    try:
        doc = fitz.open(local_pdf_path)
        output: list[str] = []

        for page_number, page in enumerate(doc, start=1):
            pix = page.get_pixmap(dpi=300)

            with tempfile.NamedTemporaryFile(suffix=".png", delete=True) as tmp:
                pix.save(tmp.name)

                page_type = classify_image_text_type(tmp.name)
                if page_type in {"handwritten_image", "mixed_image"}:
                    text = extract_handwritten_image(tmp.name)
                elif page_type == "printed_image":
                    text = extract_printed_image(tmp.name)
                else:
                    text = ""

                output.append(f"\n--- PAGE {page_number} | {page_type} ---\n{text}")

        return "\n".join(output)
    finally:
        cleanup_temporary_resource(directory=cleanup_dir)


def extract_docx(asset_path: str) -> str:
    from docx import Document

    document = Document(asset_path)
    paragraphs = [
        paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()
    ]

    tables: list[str] = []
    for table in document.tables:
        for row in table.rows:
            tables.append(" | ".join(cell.text for cell in row.cells))

    return "\n".join(paragraphs + tables)


def _chunk_dataframe_rows(dataframe: Any, rows_per_chunk: int = 200) -> list[str]:
    if dataframe.empty:
        return ["[Empty sheet]"]

    chunk_texts: list[str] = []
    total_rows = len(dataframe)

    for start in range(0, total_rows, rows_per_chunk):
        end = min(start + rows_per_chunk, total_rows)
        dataframe_chunk = dataframe.iloc[start:end]
        chunk_texts.append(
            "\n".join(
                [
                    f"--- ROWS {start + 1} TO {end} OF {total_rows} ---",
                    dataframe_chunk.to_string(index=False),
                ]
            )
        )

    return chunk_texts


def extract_excel(asset_path: str) -> str:
    import pandas as pd

    local_excel_path, cleanup_dir = resolve_remote_or_local_file(
        asset_path,
        label="Excel",
        default_suffix=".xlsx",
    )
    sheets = pd.read_excel(local_excel_path, sheet_name=None)
    output: list[str] = []

    try:
        for sheet_name, dataframe in sheets.items():
            output.append(f"\n=== SHEET: {sheet_name} ===\n")
            output.extend(_chunk_dataframe_rows(dataframe))

        return "\n".join(output)
    finally:
        cleanup_temporary_resource(directory=cleanup_dir)


def extract_pptx(asset_path: str) -> str:
    from pptx import Presentation

    presentation = Presentation(asset_path)
    output: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        output.append(f"\n--- SLIDE {slide_index} ---\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                output.append(shape.text)

    return "\n".join(output)


def get_audio_model() -> tuple[str, Any]:
    global _AUDIO_MODEL, _AUDIO_BACKEND

    if _AUDIO_MODEL is not None and _AUDIO_BACKEND is not None:
        return _AUDIO_BACKEND, _AUDIO_MODEL

    try:
        faster_whisper = importlib.import_module("faster_whisper")
        _AUDIO_MODEL = faster_whisper.WhisperModel(
            "distil-large-v3",
            device="auto",
            compute_type="int8",
        )
        _AUDIO_BACKEND = "faster_whisper"
        return _AUDIO_BACKEND, _AUDIO_MODEL
    except ModuleNotFoundError:
        pass

    try:
        whisper = importlib.import_module("whisper")
        _AUDIO_MODEL = whisper.load_model("base")
        _AUDIO_BACKEND = "openai_whisper"
        return _AUDIO_BACKEND, _AUDIO_MODEL
    except ModuleNotFoundError as exc:
        raise RuntimeError(
            "Audio transcription requires either `faster-whisper` or "
            "`openai-whisper` to be installed."
        ) from exc


def preprocess_audio(asset_path: str) -> str:
    temp_wav = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
    temp_wav.close()

    subprocess.run(
        [
            "ffmpeg",
            "-y",
            "-i",
            asset_path,
            "-ac",
            "1",
            "-ar",
            "16000",
            temp_wav.name,
        ],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        check=True,
    )

    return temp_wav.name


def _join_faster_whisper_segments(segments: Any) -> str:
    return " ".join(
        segment.text.strip() for segment in segments if segment.text.strip()
    ).strip()


def _transcribe_with_faster_whisper(model: Any, audio_path: str) -> str:
    attempts = [
        {
            "beam_size": 5,
            "vad_filter": False,
            "condition_on_previous_text": False,
            "temperature": 0.0,
            "task": "transcribe",
        },
        {
            "beam_size": 1,
            "vad_filter": False,
            "condition_on_previous_text": False,
            "temperature": 0.2,
            "task": "transcribe",
        },
    ]

    for options in attempts:
        segments, _info = model.transcribe(audio_path, **options)
        transcript = _join_faster_whisper_segments(segments)
        if transcript:
            return transcript

    return ""


def _transcribe_with_openai_whisper(model: Any, audio_path: str) -> str:
    attempts = [
        {
            "fp16": False,
            "task": "transcribe",
            "condition_on_previous_text": False,
            "temperature": 0.0,
        },
        {
            "fp16": False,
            "task": "transcribe",
            "condition_on_previous_text": False,
            "temperature": 0.2,
        },
    ]

    for options in attempts:
        result = model.transcribe(audio_path, **options)
        transcript = result.get("text", "").strip()
        if transcript:
            return transcript

    return ""


def extract_audio(asset_path: str) -> str:
    local_audio_path, cleanup_dir = resolve_remote_or_local_file(
        asset_path,
        label="Audio",
        default_suffix=".mp3",
    )
    processed_audio = preprocess_audio(local_audio_path)

    try:
        backend, model = get_audio_model()

        if backend == "faster_whisper":
            transcript = _transcribe_with_faster_whisper(model, processed_audio)
        else:
            transcript = _transcribe_with_openai_whisper(model, processed_audio)
    finally:
        cleanup_temporary_resource(file_path=processed_audio, directory=cleanup_dir)

    if not transcript:
        return (
            "[No speech could be transcribed from this audio file. "
            "The file may be instrumental, low quality, or unsupported by the "
            "current transcription backend.]"
        )

    return transcript


def download_video_url(video_url: str) -> tuple[str, str]:
    yt_dlp = importlib.import_module("yt_dlp")

    download_dir = tempfile.mkdtemp(prefix="ocr-agent-video-")
    output_template = str(Path(download_dir) / "%(title).80s.%(ext)s")

    with yt_dlp.YoutubeDL(
        {
            "outtmpl": output_template,
            "quiet": True,
            "no_warnings": True,
            "noplaylist": True,
            "format": "bestvideo*+bestaudio/best",
            "merge_output_format": "mp4",
        }
    ) as downloader:
        info = downloader.extract_info(video_url, download=True)
        requested_downloads = info.get("requested_downloads") or []
        if requested_downloads and requested_downloads[0].get("filepath"):
            return requested_downloads[0]["filepath"], download_dir

        local_path = info.get("_filename") or downloader.prepare_filename(info)
        return local_path, download_dir


def resolve_video_source(video_source: str) -> tuple[str, str | None]:
    if is_url(video_source):
        return download_video_url(video_source)

    return resolve_local_file(video_source, "Video"), None


def extract_audio_track_from_video(video_path: str) -> str:
    temp_audio = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
    temp_audio.close()

    subprocess.run(
        [
            "ffmpeg",
            "-y",
            "-i",
            video_path,
            "-vn",
            temp_audio.name,
        ],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        check=True,
    )

    return temp_audio.name


def extract_video(asset_path: str) -> str:
    resolved_video_path, cleanup_dir = resolve_video_source(asset_path)
    extracted_audio_path: str | None = None

    try:
        extracted_audio_path = extract_audio_track_from_video(resolved_video_path)
        return extract_audio(extracted_audio_path)
    finally:
        cleanup_temporary_resource(
            file_path=extracted_audio_path,
            directory=cleanup_dir,
        )


def extract_any(
    asset_path: str,
    file_type: str,
    image_text_type: str | None = None,
) -> str:
    if file_type == "text":
        return extract_text_file(asset_path)

    if file_type == "image":
        if image_text_type in {"handwritten_image", "mixed_image"}:
            return extract_handwritten_image(asset_path)
        if image_text_type == "no_text_image":
            return ""
        return extract_printed_image(asset_path)

    if file_type == "pdf":
        return extract_pdf_as_images(asset_path)

    if file_type == "docx":
        return extract_docx(asset_path)

    if file_type == "excel":
        return extract_excel(asset_path)

    if file_type == "pptx":
        return extract_pptx(asset_path)

    if file_type == "audio":
        return extract_audio(asset_path)

    if file_type == "video":
        return extract_video(asset_path)

    raise ValueError(f"Unsupported file type: {file_type}")
